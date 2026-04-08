#!/usr/bin/env python3
"""
LibreOffice Impress verification utilities for gym-anything tasks
Provides helper functions to verify Impress presentation tasks using ODP/PPTX parsing
"""

import logging
import os
import tempfile
import shutil
from pathlib import Path
from typing import Dict, List, Any, Tuple, Optional, Callable

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Import ODF modules for ODP parsing
try:
    from odf import opendocument, draw, text, presentation, style
    from odf.namespaces import DRAWNS, TEXTNS, PRESENTATIONNS
    ODF_AVAILABLE = True
except ImportError:
    logger.warning("odfpy not available - ODP parsing will be limited")
    ODF_AVAILABLE = False

# Import python-pptx for PPTX parsing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not available - PPTX parsing will be limited")
    PPTX_AVAILABLE = False


def parse_odp_file(filepath: str) -> Dict[str, Any]:
    """
    Parse an ODP (Open Document Presentation) file

    Args:
        filepath: Path to ODP file

    Returns:
        Dict containing parsed presentation data
    """
    if not ODF_AVAILABLE:
        return {'error': 'odfpy library not available'}

    try:
        doc = opendocument.load(filepath)
        
        slides_data = []
        
        # Get all draw:page elements (slides)
        for slide_elem in doc.getElementsByType(draw.Page):
            slide_name = slide_elem.getAttribute('name') or f"Slide {len(slides_data) + 1}"
            
            # Extract text content
            text_elements = []
            for text_elem in slide_elem.getElementsByType(text.P):
                text_content = str(text_elem)
                if text_content.strip():
                    text_elements.append(text_content.strip())
            
            # Check for images
            images = slide_elem.getElementsByType(draw.Image)
            has_images = len(images) > 0
            
            # Check for shapes
            shapes = []
            for shape_elem in slide_elem.getElementsByType(draw.CustomShape):
                shapes.append({
                    'type': 'custom_shape',
                    'name': shape_elem.getAttribute('name') or 'unnamed'
                })
            for shape_elem in slide_elem.getElementsByType(draw.Rect):
                shapes.append({'type': 'rectangle'})
            for shape_elem in slide_elem.getElementsByType(draw.Ellipse):
                shapes.append({'type': 'ellipse'})
            for shape_elem in slide_elem.getElementsByType(draw.Line):
                shapes.append({'type': 'line'})
            for shape_elem in slide_elem.getElementsByType(draw.Connector):
                shapes.append({'type': 'connector'})
            
            # Check for charts (usually embedded as objects)
            charts = slide_elem.getElementsByType(draw.Frame)
            has_chart = any('chart' in str(chart).lower() for chart in charts)
            
            slides_data.append({
                'name': slide_name,
                'text_elements': text_elements,
                'has_images': has_images,
                'image_count': len(images),
                'shapes': shapes,
                'shape_count': len(shapes),
                'has_chart': has_chart,
            })
        
        return {
            'filepath': filepath,
            'format': 'odp',
            'slide_count': len(slides_data),
            'slides': slides_data
        }
    
    except Exception as e:
        logger.error(f"Error parsing ODP file {filepath}: {e}")
        return {'error': str(e)}


def parse_pptx_file(filepath: str) -> Dict[str, Any]:
    """
    Parse a PPTX (PowerPoint) file

    Args:
        filepath: Path to PPTX file

    Returns:
        Dict containing parsed presentation data
    """
    if not PPTX_AVAILABLE:
        return {'error': 'python-pptx library not available'}

    try:
        prs = Presentation(filepath)
        
        slides_data = []
        
        for slide_idx, slide in enumerate(prs.slides):
            # Extract text content
            text_elements = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_elements.append(shape.text.strip())
            
            # Check for images
            images = [shape for shape in slide.shapes if shape.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
            has_images = len(images) > 0
            
            # Check for charts
            charts = [shape for shape in slide.shapes if hasattr(shape, 'chart')]
            has_chart = len(charts) > 0
            
            # Count shapes
            shapes = []
            for shape in slide.shapes:
                shapes.append({
                    'type': str(shape.shape_type),
                    'name': shape.name if hasattr(shape, 'name') else 'unnamed'
                })
            
            slides_data.append({
                'name': f"Slide {slide_idx + 1}",
                'text_elements': text_elements,
                'has_images': has_images,
                'image_count': len(images),
                'shapes': shapes,
                'shape_count': len(shapes),
                'has_chart': has_chart,
            })
        
        return {
            'filepath': filepath,
            'format': 'pptx',
            'slide_count': len(slides_data),
            'slides': slides_data
        }
    
    except Exception as e:
        logger.error(f"Error parsing PPTX file {filepath}: {e}")
        return {'error': str(e)}


def get_slide_count(data: Dict[str, Any]) -> int:
    """Get the number of slides in the presentation"""
    return data.get('slide_count', 0)


def get_slide_text_content(data: Dict[str, Any], slide_index: int) -> Tuple[Optional[str], List[str]]:
    """
    Get text content from a specific slide
    
    Args:
        data: Parsed presentation data
        slide_index: Index of slide (0-based)
        
    Returns:
        Tuple of (title, bullet_points)
    """
    slides = data.get('slides', [])
    if slide_index >= len(slides):
        return None, []
    
    slide = slides[slide_index]
    text_elements = slide.get('text_elements', [])
    
    if not text_elements:
        return None, []
    
    # First text element is usually the title
    title = text_elements[0] if text_elements else None
    bullets = text_elements[1:] if len(text_elements) > 1 else []
    
    return title, bullets


def get_slide_title(data: Dict[str, Any], slide_index: int) -> Optional[str]:
    """Get the title of a specific slide"""
    title, _ = get_slide_text_content(data, slide_index)
    return title


def get_slide_bullets(data: Dict[str, Any], slide_index: int) -> List[str]:
    """Get bullet points from a specific slide"""
    _, bullets = get_slide_text_content(data, slide_index)
    return bullets


def check_slide_has_images(data: Dict[str, Any], slide_index: int) -> bool:
    """Check if a slide contains images"""
    slides = data.get('slides', [])
    if slide_index >= len(slides):
        return False
    
    return slides[slide_index].get('has_images', False)


def check_slide_has_chart(data: Dict[str, Any], slide_index: int) -> bool:
    """Check if a slide contains a chart"""
    slides = data.get('slides', [])
    if slide_index >= len(slides):
        return False
    
    return slides[slide_index].get('has_chart', False)


def check_slide_has_shapes(data: Dict[str, Any], slide_index: int, min_count: int = 1) -> bool:
    """Check if a slide contains shapes"""
    slides = data.get('slides', [])
    if slide_index >= len(slides):
        return False
    
    shape_count = slides[slide_index].get('shape_count', 0)
    return shape_count >= min_count


def check_slide_has_animations(data: Dict[str, Any], slide_index: int) -> bool:
    """
    Check if a slide has animations
    
    Note: Full animation parsing from ODP/PPTX is complex.
    This is a simplified check that may require file inspection.
    """
    # This would require more sophisticated parsing of presentation:animations
    # For now, return a placeholder
    logger.warning("Animation detection not fully implemented")
    return False


def verify_slide_transition(data: Dict[str, Any], slide_index: int) -> bool:
    """
    Check if a slide has a transition effect
    
    Note: Full transition parsing from ODP/PPTX is complex.
    This is a simplified check.
    """
    # This would require parsing presentation:transition elements
    logger.warning("Transition detection not fully implemented")
    return False


def get_slide_shapes(data: Dict[str, Any], slide_index: int) -> List[Dict[str, str]]:
    """Get list of shapes on a specific slide"""
    slides = data.get('slides', [])
    if slide_index >= len(slides):
        return []
    
    return slides[slide_index].get('shapes', [])


def verify_text_on_slide(data: Dict[str, Any], slide_index: int, expected_text: str, 
                        case_sensitive: bool = False) -> bool:
    """
    Verify that specific text appears on a slide
    
    Args:
        data: Parsed presentation data
        slide_index: Index of slide
        expected_text: Text to search for
        case_sensitive: Whether to perform case-sensitive search
        
    Returns:
        True if text is found, False otherwise
    """
    slides = data.get('slides', [])
    if slide_index >= len(slides):
        return False
    
    text_elements = slides[slide_index].get('text_elements', [])
    all_text = ' '.join(text_elements)
    
    if not case_sensitive:
        all_text = all_text.lower()
        expected_text = expected_text.lower()
    
    return expected_text in all_text


def count_shapes_on_slide(data: Dict[str, Any], slide_index: int, shape_type: Optional[str] = None) -> int:
    """
    Count shapes on a slide, optionally filtered by type
    
    Args:
        data: Parsed presentation data
        slide_index: Index of slide
        shape_type: Optional shape type to filter by
        
    Returns:
        Count of shapes
    """
    shapes = get_slide_shapes(data, slide_index)
    
    if shape_type is None:
        return len(shapes)
    
    return sum(1 for shape in shapes if shape.get('type') == shape_type)


def get_presentation_metadata(data: Dict[str, Any]) -> Dict[str, Any]:
    """Get metadata about the presentation"""
    return {
        'format': data.get('format'),
        'slide_count': data.get('slide_count', 0),
        'filepath': data.get('filepath'),
    }


def setup_verification_environment(copy_from_env_fn: Callable,
                                  container_path: str,
                                  expected_formats: List[str] = None) -> Tuple[bool, Dict[str, Any]]:
    """
    Set up verification environment by copying presentation file from container

    Args:
        copy_from_env_fn: Function to copy files from container
        container_path: Path to presentation file in container
        expected_formats: List of expected formats (['odp', 'pptx'])

    Returns:
        Tuple of (success, data_dict)
        data_dict contains: {'filepath': str, 'data': parsed_data, 'temp_dir': str}
    """
    if expected_formats is None:
        expected_formats = ['odp', 'pptx']

    # Create unique temporary directory
    temp_dir = Path(tempfile.mkdtemp(prefix='impress_verify_'))

    try:
        # Determine file extension
        file_ext = Path(container_path).suffix.lower()
        host_file = temp_dir / f"result{file_ext}"

        # Copy file from container
        try:
            copy_from_env_fn(container_path, str(host_file))
        except Exception as e:
            shutil.rmtree(temp_dir, ignore_errors=True)
            return False, {'error': f"Failed to copy file: {e}"}

        if not host_file.exists() or host_file.stat().st_size == 0:
            shutil.rmtree(temp_dir, ignore_errors=True)
            return False, {'error': f"File not found or empty: {container_path}"}

        # Parse file based on format
        if file_ext == '.odp' and 'odp' in expected_formats:
            data = parse_odp_file(str(host_file))
        elif file_ext in ['.pptx', '.ppt'] and 'pptx' in expected_formats:
            data = parse_pptx_file(str(host_file))
        else:
            shutil.rmtree(temp_dir, ignore_errors=True)
            return False, {'error': f"Unsupported file format: {file_ext}"}

        if 'error' in data:
            shutil.rmtree(temp_dir, ignore_errors=True)
            return False, {'error': f"Parse error: {data['error']}"}

        return True, {
            'filepath': str(host_file),
            'data': data,
            'temp_dir': str(temp_dir)
        }

    except Exception as e:
        shutil.rmtree(temp_dir, ignore_errors=True)
        logger.error(f"Setup verification environment failed: {e}")
        return False, {'error': str(e)}


def cleanup_verification_environment(temp_dir: Optional[str] = None):
    """
    Clean up temporary verification files

    Args:
        temp_dir: Path to temp directory to clean up
    """
    if temp_dir is None:
        logger.warning("cleanup_verification_environment called with temp_dir=None, skipping cleanup")
        return

    if os.path.exists(temp_dir):
        try:
            shutil.rmtree(temp_dir)
            logger.debug(f"Cleaned up temp directory: {temp_dir}")
        except Exception as e:
            logger.warning(f"Failed to cleanup temp directory {temp_dir}: {e}")


def copy_and_parse_presentation(container_path: str, copy_from_env_fn: Callable,
                               file_format: str = 'odp') -> Tuple[bool, Dict[str, Any], str, str]:
    """
    Copy presentation from container and parse it.

    Args:
        container_path: Path to file in container
        copy_from_env_fn: Function to copy files from container
        file_format: Expected file format ('odp' or 'pptx')

    Returns:
        Tuple of (success, parsed_data, error_message, temp_dir)
    """
    success, result = setup_verification_environment(copy_from_env_fn, container_path, [file_format])

    if not success:
        return False, {}, result.get('error', 'Unknown error'), ''

    return True, result.get('data', {}), "", result.get('temp_dir', '')
