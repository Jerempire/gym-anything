#!/usr/bin/env python3
"""
Verifier for insert_data_table task.
Checks if the agent correctly inserted and populated a table in LibreOffice Impress.
"""

import json
import tempfile
import os
import logging
import sys
import shutil

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Try to import presentation parsing libraries
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from odf import opendocument, table, text, draw
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False


def clean_text(text):
    """Normalize text for comparison."""
    if not text:
        return ""
    return str(text).strip().lower()


def extract_table_data_pptx(filepath, slide_index):
    """Extract table data from a PPTX file."""
    if not PPTX_AVAILABLE:
        return None, "python-pptx library not available"
    
    try:
        prs = Presentation(filepath)
        if slide_index >= len(prs.slides):
            return None, f"Slide index {slide_index} out of range (count: {len(prs.slides)})"
        
        slide = prs.slides[slide_index]
        tables = [shape.table for shape in slide.shapes if shape.has_table]
        
        if not tables:
            return None, "No table found on target slide"
        
        # Assume the user worked on the first/primary table
        tbl = tables[0]
        data = []
        for row in tbl.rows:
            row_data = [cell.text_frame.text for cell in row.cells]
            data.append(row_data)
            
        return data, None
    except Exception as e:
        return None, str(e)


def extract_table_data_odp(filepath, slide_index):
    """Extract table data from an ODP file."""
    if not ODF_AVAILABLE:
        return None, "odfpy library not available"
    
    try:
        doc = opendocument.load(filepath)
        slides = doc.getElementsByType(draw.Page)
        
        if slide_index >= len(slides):
            return None, f"Slide index {slide_index} out of range (count: {len(slides)})"
        
        slide = slides[slide_index]
        tables = slide.getElementsByType(table.Table)
        
        if not tables:
            # Tables might be inside frames
            frames = slide.getElementsByType(draw.Frame)
            for frame in frames:
                tables.extend(frame.getElementsByType(table.Table))
        
        if not tables:
            return None, "No table found on target slide"
        
        tbl = tables[0]
        data = []
        rows = tbl.getElementsByType(table.TableRow)
        
        for row in rows:
            row_data = []
            cells = row.getElementsByType(table.TableCell)
            for cell in cells:
                # Extract text from paragraphs within cell
                paras = cell.getElementsByType(text.P)
                cell_text = "\n".join([str(p) for p in paras])
                row_data.append(cell_text)
            data.append(row_data)
            
        return data, None
    except Exception as e:
        return None, str(e)


def verify_table_data(traj, env_info, task_info):
    """
    Verify the inserted table data.
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    # 1. Load result JSON
    temp_result = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", temp_result.name)
        with open(temp_result.name, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to load result JSON: {e}"}
    finally:
        if os.path.exists(temp_result.name):
            os.unlink(temp_result.name)

    # Basic Checks
    if not result.get('file_exists'):
        return {"passed": False, "score": 0, "feedback": "Target presentation file not found"}
    
    if not result.get('file_modified_during_task'):
        return {"passed": False, "score": 0, "feedback": "Presentation file was not modified/saved during the task"}

    # 2. Retrieve Presentation File
    file_path = result.get('file_path')
    file_format = result.get('file_format')
    
    temp_pres = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_format}')
    try:
        copy_from_env(file_path, temp_pres.name)
        
        # 3. Parse Table Data
        slide_index = task_info['metadata'].get('slide_index', 2)
        
        if file_format == 'pptx':
            table_data, error = extract_table_data_pptx(temp_pres.name, slide_index)
        elif file_format == 'odp':
            table_data, error = extract_table_data_odp(temp_pres.name, slide_index)
        else:
            return {"passed": False, "score": 0, "feedback": f"Unknown file format: {file_format}"}
            
        if error:
            return {"passed": False, "score": 20, "feedback": f"File modified but table parsing failed: {error}"}
            
        if not table_data:
            return {"passed": False, "score": 20, "feedback": "File modified but no data found in table"}

    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Error retrieving/parsing file: {e}"}
    finally:
        if os.path.exists(temp_pres.name):
            os.unlink(temp_pres.name)

    # 4. Verify Content
    score = 0
    feedback_lines = []
    
    # Check dimensions
    num_rows = len(table_data)
    num_cols = len(table_data[0]) if num_rows > 0 else 0
    
    expected_rows = task_info['metadata'].get('expected_rows', 6)
    expected_cols = task_info['metadata'].get('expected_cols', 4)
    
    if num_rows >= expected_rows and num_cols >= expected_cols:
        score += 20
        feedback_lines.append(f"✅ Table dimensions correct ({num_rows}x{num_cols})")
    else:
        feedback_lines.append(f"❌ Incorrect table dimensions: found {num_rows}x{num_cols}, expected {expected_rows}x{expected_cols}")
        # Return early if table is empty or too small to evaluate
        if num_rows < 2:
            return {"passed": False, "score": score, "feedback": " | ".join(feedback_lines)}

    # Check Headers
    expected_headers = task_info['metadata'].get('expected_headers', [])
    header_matches = 0
    if num_rows > 0:
        row0 = [clean_text(c) for c in table_data[0]]
        for exp in expected_headers:
            if any(clean_text(exp) in cell for cell in row0):
                header_matches += 1
    
    if header_matches >= 3:
        score += 20
        feedback_lines.append(f"✅ Headers correct ({header_matches}/4)")
    else:
        feedback_lines.append(f"❌ Headers missing or incorrect (found {header_matches}/4)")

    # Check Data Rows
    # We look for key values anywhere in the table to be lenient about exact row ordering
    # providing the country name is present in the row
    expected_data = task_info['metadata'].get('expected_data', [])
    data_points_score = 0
    max_data_points = len(expected_data) * 4 # 4 fields per row
    
    correct_rows = 0
    
    for exp_row in expected_data:
        country = clean_text(exp_row['country'])
        capacity = str(exp_row['capacity'])
        growth = str(exp_row['growth'])
        share = str(exp_row['share'])
        
        row_found = False
        for table_row in table_data[1:]: # Skip header
            row_text = " ".join([clean_text(c) for c in table_row])
            
            # If this row contains the country name
            if country in row_text:
                row_found = True
                row_score = 0
                
                # Check other values
                # We use simple substring matching for robustness against formatting (e.g. "393" vs "393.0")
                if capacity in row_text: row_score += 1
                if growth in row_text: row_score += 1
                if share in row_text: row_score += 1
                
                # Add score for finding the country itself
                row_score += 1 
                
                data_points_score += row_score
                
                if row_score >= 3: # Mostly correct row
                    correct_rows += 1
                break
    
    # Normalize data score to 60 points max
    normalized_data_score = int((data_points_score / max_data_points) * 60)
    score += normalized_data_score
    
    feedback_lines.append(f"Data accuracy score: {normalized_data_score}/60")
    if correct_rows >= 4:
         feedback_lines.append(f"✅ At least 4/5 rows contain correct data")
    else:
         feedback_lines.append(f"⚠️ Only {correct_rows}/5 rows fully correct")

    return {
        "passed": score >= 70,
        "score": score,
        "feedback": " | ".join(feedback_lines)
    }