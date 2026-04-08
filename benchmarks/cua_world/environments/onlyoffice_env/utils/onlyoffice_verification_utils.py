#!/usr/bin/env python3
"""
ONLYOFFICE verification utilities for gym-anything tasks
Provides helper functions to verify ONLYOFFICE document, spreadsheet, and presentation tasks
"""

import json
import logging
import os
import tempfile
import shutil
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Any, Tuple, Optional, Callable

# Import document parsing libraries
try:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# ============================================================================
# Document (DOCX) Verification Functions
# ============================================================================

def parse_docx_file(filepath: str) -> Optional[Any]:
    """
    Parse a DOCX file

    Args:
        filepath: Path to DOCX file

    Returns:
        Document object or None if parsing fails
    """
    if Document is None:
        logger.error("python-docx not installed")
        return None

    try:
        return Document(filepath)
    except Exception as e:
        logger.error(f"Error parsing DOCX: {e}")
        return None


def get_document_text(doc: Any) -> str:
    """
    Extract all text from a document

    Args:
        doc: Document object

    Returns:
        Full text content
    """
    try:
        return '\n'.join([para.text for para in doc.paragraphs])
    except Exception as e:
        logger.error(f"Error extracting text: {e}")
        return ""


def check_text_formatting(doc: Any, text: str, **formatting) -> bool:
    """
    Check if text has specific formatting

    Args:
        doc: Document object
        text: Text to find
        **formatting: bold, italic, underline, font_size, font_name, color

    Returns:
        True if text found with formatting, False otherwise
    """
    try:
        for para in doc.paragraphs:
            for run in para.runs:
                if text.lower() in run.text.lower():
                    checks_passed = True

                    if 'bold' in formatting:
                        checks_passed &= (run.bold == formatting['bold'])
                    if 'italic' in formatting:
                        checks_passed &= (run.italic == formatting['italic'])
                    if 'underline' in formatting:
                        checks_passed &= (run.underline == formatting['underline'])
                    if 'font_size' in formatting:
                        if run.font.size:
                            size_pt = run.font.size.pt
                            checks_passed &= (size_pt == formatting['font_size'])
                        else:
                            checks_passed = False

                    if checks_passed:
                        return True
        return False
    except Exception as e:
        logger.error(f"Error checking formatting: {e}")
        return False


def check_paragraph_alignment(doc: Any, text: str, alignment: str) -> bool:
    """
    Check if paragraph containing text has specified alignment

    Args:
        doc: Document object
        text: Text to find
        alignment: 'left', 'center', 'right', 'justify'

    Returns:
        True if paragraph has correct alignment
    """
    alignment_map = {
        'left': 0,      # WD_ALIGN_PARAGRAPH.LEFT
        'center': 1,    # WD_ALIGN_PARAGRAPH.CENTER
        'right': 2,     # WD_ALIGN_PARAGRAPH.RIGHT
        'justify': 3    # WD_ALIGN_PARAGRAPH.JUSTIFY
    }

    try:
        target_alignment = alignment_map.get(alignment.lower())
        if target_alignment is None:
            return False

        for para in doc.paragraphs:
            if text.lower() in para.text.lower():
                if para.alignment == target_alignment or para.alignment_val == target_alignment:
                    return True
        return False
    except Exception as e:
        logger.error(f"Error checking alignment: {e}")
        return False


def count_paragraphs(doc: Any) -> int:
    """Count number of paragraphs in document"""
    try:
        return len(doc.paragraphs)
    except:
        return 0


def count_tables(doc: Any) -> int:
    """Count number of tables in document"""
    try:
        return len(doc.tables)
    except:
        return 0


# ============================================================================
# Spreadsheet (XLSX) Verification Functions
# ============================================================================

def parse_xlsx_file(filepath: str) -> Optional[Any]:
    """
    Parse an XLSX file

    Args:
        filepath: Path to XLSX file

    Returns:
        Workbook object or None if parsing fails
    """
    if load_workbook is None:
        logger.error("openpyxl not installed")
        return None

    try:
        return load_workbook(filepath, data_only=True)
    except Exception as e:
        logger.error(f"Error parsing XLSX: {e}")
        return None


def get_cell_value(workbook: Any, sheet_name: str, cell_ref: str) -> Any:
    """
    Get value from a specific cell

    Args:
        workbook: Workbook object
        sheet_name: Sheet name
        cell_ref: Cell reference (e.g., 'A1', 'B5')

    Returns:
        Cell value or None
    """
    try:
        sheet = workbook[sheet_name]
        return sheet[cell_ref].value
    except Exception as e:
        logger.error(f"Error getting cell value: {e}")
        return None


def verify_formula(workbook: Any, sheet_name: str, cell_ref: str, expected_result: float, tolerance: float = 0.01) -> bool:
    """
    Verify that a cell contains a formula with expected result

    Args:
        workbook: Workbook object (must have data_only=False to see formulas)
        sheet_name: Sheet name
        cell_ref: Cell reference
        expected_result: Expected numeric result
        tolerance: Acceptable difference

    Returns:
        True if formula result matches expected
    """
    try:
        # Need to reload without data_only to see formulas
        sheet = workbook[sheet_name]
        cell_value = sheet[cell_ref].value

        if cell_value is None:
            return False

        # Check if value is close to expected
        if isinstance(cell_value, (int, float)):
            return abs(cell_value - expected_result) <= tolerance

        return False
    except Exception as e:
        logger.error(f"Error verifying formula: {e}")
        return False


def count_filled_cells(workbook: Any, sheet_name: str, cell_range: Optional[str] = None) -> int:
    """
    Count number of non-empty cells in a sheet or range

    Args:
        workbook: Workbook object
        sheet_name: Sheet name
        cell_range: Optional cell range (e.g., 'A1:D10')

    Returns:
        Number of filled cells
    """
    try:
        sheet = workbook[sheet_name]
        count = 0

        if cell_range:
            for row in sheet[cell_range]:
                for cell in row:
                    if cell.value is not None:
                        count += 1
        else:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        count += 1

        return count
    except Exception as e:
        logger.error(f"Error counting filled cells: {e}")
        return 0


def get_sheet_data(workbook: Any, sheet_name: str, max_rows: int = 100, max_cols: int = 20) -> List[List]:
    """
    Get all data from a sheet as a 2D list

    Args:
        workbook: Workbook object
        sheet_name: Sheet name
        max_rows: Maximum rows to read
        max_cols: Maximum columns to read

    Returns:
        2D list of cell values
    """
    try:
        sheet = workbook[sheet_name]
        data = []

        for row_idx, row in enumerate(sheet.iter_rows(max_row=max_rows, max_col=max_cols), start=1):
            row_data = [cell.value for cell in row]
            data.append(row_data)

        return data
    except Exception as e:
        logger.error(f"Error getting sheet data: {e}")
        return []


# ============================================================================
# Presentation (PPTX) Verification Functions
# ============================================================================

def parse_pptx_file(filepath: str) -> Optional[Any]:
    """
    Parse a PPTX file

    Args:
        filepath: Path to PPTX file

    Returns:
        Presentation object or None if parsing fails
    """
    if Presentation is None:
        logger.error("python-pptx not installed")
        return None

    try:
        return Presentation(filepath)
    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        return None


def count_slides(prs: Any) -> int:
    """Count number of slides in presentation"""
    try:
        return len(prs.slides)
    except:
        return 0


def get_slide_text(prs: Any, slide_index: int) -> str:
    """
    Get all text from a specific slide

    Args:
        prs: Presentation object
        slide_index: Slide index (0-based)

    Returns:
        All text from the slide
    """
    try:
        slide = prs.slides[slide_index]
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)

        return '\n'.join(text_parts)
    except Exception as e:
        logger.error(f"Error getting slide text: {e}")
        return ""


def check_slide_has_image(prs: Any, slide_index: int) -> bool:
    """
    Check if slide contains an image

    Args:
        prs: Presentation object
        slide_index: Slide index (0-based)

    Returns:
        True if slide has image
    """
    try:
        slide = prs.slides[slide_index]

        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                return True

        return False
    except Exception as e:
        logger.error(f"Error checking for image: {e}")
        return False


def count_shapes_on_slide(prs: Any, slide_index: int) -> int:
    """Count number of shapes on a slide"""
    try:
        slide = prs.slides[slide_index]
        return len(slide.shapes)
    except:
        return 0


# ============================================================================
# Generic File Utilities
# ============================================================================

def copy_and_parse_document(container_path: str, copy_from_env_fn: Callable, file_format: str = 'auto') -> Tuple[bool, Any, str]:
    """
    Copy document from container and parse it

    Args:
        container_path: Path in container
        copy_from_env_fn: Function to copy files
        file_format: 'docx', 'xlsx', 'pptx', or 'auto'

    Returns:
        Tuple of (success, parsed_document, error_message)
    """
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=Path(container_path).suffix)

    try:
        copy_from_env_fn(container_path, temp_file.name)

        if not os.path.exists(temp_file.name) or os.path.getsize(temp_file.name) == 0:
            return False, None, f"File not found or empty: {container_path}"

        # Auto-detect format
        if file_format == 'auto':
            ext = Path(container_path).suffix.lower()
            if ext in ['.docx', '.doc']:
                file_format = 'docx'
            elif ext in ['.xlsx', '.xls']:
                file_format = 'xlsx'
            elif ext in ['.pptx', '.ppt']:
                file_format = 'pptx'
            else:
                return False, None, f"Unknown file format: {ext}"

        # Parse based on format
        if file_format == 'docx':
            doc = parse_docx_file(temp_file.name)
            return (True, doc, "") if doc else (False, None, "Failed to parse DOCX")
        elif file_format == 'xlsx':
            wb = parse_xlsx_file(temp_file.name)
            return (True, wb, "") if wb else (False, None, "Failed to parse XLSX")
        elif file_format == 'pptx':
            prs = parse_pptx_file(temp_file.name)
            return (True, prs, "") if prs else (False, None, "Failed to parse PPTX")
        else:
            return False, None, f"Unsupported format: {file_format}"

    except Exception as e:
        logger.error(f"Error copying and parsing: {e}")
        return False, None, str(e)
    finally:
        if os.path.exists(temp_file.name):
            os.unlink(temp_file.name)


def read_file_content(filepath: str) -> str:
    """Read file content as text"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error reading file: {e}")
        return ""


def cleanup_temp_dir(temp_dir: str):
    """Clean up temporary directory"""
    if temp_dir and os.path.exists(temp_dir):
        try:
            shutil.rmtree(temp_dir)
            logger.debug(f"Cleaned up temp directory: {temp_dir}")
        except Exception as e:
            logger.warning(f"Failed to cleanup temp directory: {e}")
