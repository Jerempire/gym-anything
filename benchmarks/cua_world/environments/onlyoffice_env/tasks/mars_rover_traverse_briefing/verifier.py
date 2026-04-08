#!/usr/bin/env python3
"""
Verifier for Mars Rover Traverse Briefing task.

Requirements Verified Programmatically:
1. File exists and was created during the task run (10 pts)
2. Presentation has at least 4 slides (10 pts)
3. Text content covers Key Objectives (20 pts)
4. Image shape successfully added (20 pts)
5. Table successfully added (10 pts)
6. Data Filtering Accuracy: High-priority data is present, Low/Medium distractors are missing (30 pts)
"""

import sys
import os
import json
import logging
import tempfile

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Check for python-pptx dependency gracefully
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. Will try to install.")

def ensure_dependencies():
    global PPTX_AVAILABLE
    if not PPTX_AVAILABLE:
        import subprocess
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            PPTX_AVAILABLE = True
        except Exception as e:
            logger.error(f"Failed to install python-pptx: {e}")
            return False
    return True

def verify_mars_briefing(traj, env_info, task_info):
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0.0, "feedback": "Copy function not available"}

    if not ensure_dependencies():
        return {"passed": False, "score": 0.0, "feedback": "Failed to import or install python-pptx requirement"}

    from pptx import Presentation
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        PICTURE_TYPE = MSO_SHAPE_TYPE.PICTURE
    except:
        PICTURE_TYPE = 13  # Fallback constant for Picture

    score = 0.0
    feedback = []

    # 1. Read JSON exported by environment
    temp_json = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/mars_task_result.json", temp_json.name)
        with open(temp_json.name, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0.0, "feedback": f"Failed to read result JSON: {e}"}
    finally:
        if os.path.exists(temp_json.name):
            os.unlink(temp_json.name)

    # 2. Gate check: Did they create the file during the task?
    if not result.get("output_file_exists", False):
        return {"passed": False, "score": 0.0, "feedback": "Presentation file not found at expected path."}

    if not result.get("file_created_during_task", False):
        return {"passed": False, "score": 0.0, "feedback": "Presentation file was not created or modified during the task."}

    score += 10.0
    feedback.append("File creation check passed.")

    # 3. Parse PPTX via temporary file
    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    try:
        copy_from_env("/home/ga/Documents/Presentations/traverse_briefing.pptx", temp_pptx.name)
        prs = Presentation(temp_pptx.name)
    except Exception as e:
        return {"passed": False, "score": score, "feedback": f"Failed to parse PPTX file: {e}"}

    # 4. Check Slides Count
    num_slides = len(prs.slides)
    if num_slides >= 4:
        score += 10.0
        feedback.append(f"Slide count check passed ({num_slides} slides).")
    else:
        feedback.append(f"Slide count check failed (Found {num_slides}, expected >= 4).")

    # 5. Extract Contents
    all_text = ""
    has_picture = False
    has_table = False
    table_text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text + " "
            
            # Picture check
            if getattr(shape, "shape_type", None) == PICTURE_TYPE or getattr(shape, "shape_type", None) == 13:
                has_picture = True
            
            # Table check
            if shape.has_table:
                has_table = True
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            table_text += cell.text + " "
                            all_text += cell.text + " "

    all_text = all_text.lower()
    table_text = table_text.lower()

    # 6. Verify Text Items
    expected_terms = ["jezero", "traverse", "stratigraphy", "phyllosilicate", "habitability", "sample"]
    found_terms = [t for t in expected_terms if t in all_text]
    if len(found_terms) >= 4:
        score += 20.0
        feedback.append("Text content check passed (Found key briefing notes).")
    else:
        feedback.append(f"Text content check failed (Found {len(found_terms)}/6 key terms).")

    # 7. Verify Image
    if has_picture:
        score += 20.0
        feedback.append("Image check passed (Found picture shape).")
    else:
        feedback.append("Image check failed (No picture shape found).")

    # 8. Verify Table Presence
    if has_table:
        score += 10.0
        feedback.append("Table shape check passed.")
    else:
        feedback.append("Table shape check failed.")

    # 9. Verify Table Data Filtering
    # High Priority Targets that MUST be present
    high_targets = ["rochette", "brac", "quartier", "artuby"]
    found_high = [t for t in high_targets if t in table_text]

    # Distractor Targets that MUST NOT be present (tests sorting/filtering skills)
    distractors = ["nav_point", "drive_stop", "hazcam", "uhf", "citadelle"]
    found_distractors = [t for t in distractors if t in table_text]

    if has_table:
        if len(found_high) == 4 and len(found_distractors) == 0:
            score += 30.0
            feedback.append("Table data check passed perfectly (All high priority, no distractors).")
        elif len(found_high) > 0:
            # Partial credit logic
            pts = len(found_high) * 7.5
            penalty = len(found_distractors) * 5.0
            pts = max(0, pts - penalty)
            score += pts
            feedback.append(f"Table data check partial (Found {len(found_high)} high targets, {len(found_distractors)} distractors. +{pts:.1f} pts).")
        else:
            feedback.append("Table data check failed (High priority targets missing).")

    # Evaluate Pass/Fail status
    passed = (score >= 70.0)

    # Cleanup
    if os.path.exists(temp_pptx.name):
        os.unlink(temp_pptx.name)

    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback)
    }