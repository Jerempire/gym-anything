#!/usr/bin/env python3
"""
Verifier for Mount Rainier Lahar Hazard Briefing Task.
Extracts and parses the final PPTX to check for 5 slides, table insertion,
image insertion, and required text.
"""

import os
import json
import logging
import tempfile

# Attempt to import python-pptx
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Fallback to importing VLM if PPTX parsing fails (or for hybrid verification)
import sys
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), '../../', 'utils'))

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def extract_text_from_slide(slide):
    """Extract all text from all shapes in a slide."""
    text_runs = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_runs.append(shape.text.lower())
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text_runs.append(cell.text.lower())
    return " ".join(text_runs)


def has_table(slide):
    """Check if slide contains a table."""
    for shape in slide.shapes:
        if shape.has_table:
            return True
    return False


def has_picture(slide):
    """Check if slide contains a picture."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    return False


def verify_volcano_hazard_briefing(traj, env_info, task_info):
    """
    Verify the PPTX generation task.
    
    Criteria (Max 100 points):
    - File exists & created during task (15 points) -> Gatekeeper
    - Total slides == 5 (15 points)
    - Slide 3 contains Table and "Osceola" (25 points)
    - Slide 4 contains "80,000" / "population" (20 points)
    - Slide 5 contains an Image and "evacuation" (25 points)
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    # Load result metadata
    result_json_path = tempfile.NamedTemporaryFile(delete=False, suffix='.json').name
    try:
        copy_from_env("/tmp/task_result.json", result_json_path)
        with open(result_json_path, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to read result JSON: {e}"}
    finally:
        if os.path.exists(result_json_path):
            os.unlink(result_json_path)

    score = 0
    feedback_parts = []
    
    # 1. Gatekeeper: File exists
    if not result.get('output_exists', False):
        return {
            "passed": False, 
            "score": 0, 
            "feedback": "Final presentation file (Rainier_Hazard_Briefing_Final.pptx) not found."
        }
        
    if result.get('file_created_during_task', False):
        score += 15
        feedback_parts.append("File created during task (+15)")
    else:
        feedback_parts.append("Warning: File timestamp indicates it may not have been created during this session.")

    if not PPTX_AVAILABLE:
        # If the environment lacks python-pptx for some reason, we'd fall back, 
        # but the container is guaranteed to have python-pptx per env.json.
        return {"passed": False, "score": score, "feedback": "python-pptx not available for verification."}

    # 2. Extract and Parse PPTX
    pptx_temp_path = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx').name
    try:
        copy_from_env("/home/ga/Documents/Presentations/Rainier_Hazard_Briefing_Final.pptx", pptx_temp_path)
        try:
            prs = Presentation(pptx_temp_path)
        except Exception as e:
            return {
                "passed": False, 
                "score": score, 
                "feedback": f"Failed to parse PPTX file (may be corrupted): {e}"
            }
            
        slides = prs.slides
        num_slides = len(slides)
        
        # 3. Slide Count
        if num_slides == 5:
            score += 15
            feedback_parts.append("Slide count is exactly 5 (+15)")
        else:
            feedback_parts.append(f"Expected 5 slides, found {num_slides}")
            
        # We only check slides 3, 4, 5 if they exist
        # If they added slides but maybe in wrong order, we will check index 2, 3, 4 (0-based)
        if num_slides >= 3:
            slide3 = slides[2]
            s3_text = extract_text_from_slide(slide3)
            s3_score = 0
            if has_table(slide3):
                s3_score += 15
            if "osceola" in s3_text:
                s3_score += 10
            
            score += s3_score
            if s3_score == 25:
                feedback_parts.append("Slide 3 content correct (+25)")
            else:
                feedback_parts.append(f"Slide 3 missing required elements (Score: {s3_score}/25)")
                
        if num_slides >= 4:
            slide4 = slides[3]
            s4_text = extract_text_from_slide(slide4)
            s4_score = 0
            if "80,000" in s4_text or "80000" in s4_text:
                s4_score += 15
            if "population" in s4_text:
                s4_score += 5
                
            score += s4_score
            if s4_score == 20:
                feedback_parts.append("Slide 4 content correct (+20)")
            else:
                feedback_parts.append(f"Slide 4 missing required elements (Score: {s4_score}/20)")
                
        if num_slides >= 5:
            slide5 = slides[4]
            s5_text = extract_text_from_slide(slide5)
            s5_score = 0
            if has_picture(slide5):
                s5_score += 15
            if "evacuation" in s5_text:
                s5_score += 10
                
            score += s5_score
            if s5_score == 25:
                feedback_parts.append("Slide 5 content correct (+25)")
            else:
                feedback_parts.append(f"Slide 5 missing required elements (Score: {s5_score}/25)")

    finally:
        if os.path.exists(pptx_temp_path):
            os.unlink(pptx_temp_path)

    passed = score >= 75
    
    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }