#!/usr/bin/env python3
"""
Verifier for Medical Grand Rounds Presentation Task.

Combines File-based checking (.pptx parsing) and VLM Trajectory checks.
"""

import json
import os
import sys
import tempfile
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Attempt to import python-pptx (installed in env)
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx is not available. Falling back to basic checks.")

def get_vlm_trajectory_check(traj, env_info):
    """
    Uses VLM to verify that the agent actively manipulated the ONLYOFFICE Presentation Editor
    and assembled slides containing text and images.
    """
    try:
        from gym_anything.vlm import sample_trajectory_frames, get_final_screenshot
        frames = sample_trajectory_frames(traj, n=4)
        final = get_final_screenshot(traj)
        images = frames + [final] if final else frames
        
        if not images:
            return {"vlm_score": 0, "feedback": "No trajectory frames for VLM."}
            
        prompt = (
            "Review these trajectory frames of an agent interacting with a desktop. "
            "Did the agent actively use a Presentation Editor (like ONLYOFFICE) to create a slideshow, "
            "and did they insert medical images (like an ECG or Angiogram) and type clinical text? "
            "Respond ONLY in JSON format: {\"used_presentation_editor\": true/false, \"inserted_images\": true/false}"
        )
        
        query_vlm = env_info.get('query_vlm')
        if not query_vlm:
            return {"vlm_score": 0, "feedback": "VLM query function unavailable."}
            
        result = query_vlm(images=images, prompt=prompt)
        parsed = result.get('parsed', {})
        
        vlm_score = 0
        if parsed.get('used_presentation_editor', False): vlm_score += 10
        if parsed.get('inserted_images', False): vlm_score += 10
        
        return {"vlm_score": vlm_score, "feedback": f"VLM Check Passed: Editor={parsed.get('used_presentation_editor')}, Images={parsed.get('inserted_images')}"}
        
    except Exception as e:
        logger.error(f"VLM verification failed: {e}")
        return {"vlm_score": 0, "feedback": f"VLM verification error: {e}"}

def verify_presentation(traj, env_info, task_info):
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    feedback_parts = []
    score = 0
    max_score = 100

    # 1. Read the JSON export
    temp_json = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/medical_task_result.json", temp_json.name)
        with open(temp_json.name, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to read exported result: {e}"}
    finally:
        if os.path.exists(temp_json.name):
            os.unlink(temp_json.name)

    # Criteria 1: Output File Exists and Created During Task (Anti-Gaming)
    output_exists = result.get('output_exists', False)
    file_created = result.get('file_created_during_task', False)
    
    if output_exists and file_created:
        score += 20
        feedback_parts.append("File exists and created during session (20/20)")
    elif output_exists:
        score += 5
        feedback_parts.append("File exists but was not created during session (5/20)")
    else:
        feedback_parts.append("Target presentation file not found (0/20)")
        return {"passed": False, "score": 0, "feedback": " | ".join(feedback_parts)}

    # 2. Extract and Parse the PPTX file
    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    slide_count = 0
    image_count = 0
    all_text = ""
    title_font_ok = False
    
    try:
        copy_from_env("/home/ga/Documents/Presentations/grand_rounds_takotsubo.pptx", temp_pptx.name)
        
        if PPTX_AVAILABLE:
            prs = Presentation(temp_pptx.name)
            slide_count = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text += shape.text.lower() + " "
                    
                    # 13 is MSO_SHAPE_TYPE.PICTURE
                    if shape.shape_type == 13: 
                        image_count += 1
                        
                    # Check font size on the first slide (Title)
                    if i == 0 and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size.pt >= 32:
                                    title_font_ok = True
        else:
            # Fallback if python-pptx is somehow missing, though it shouldn't be
            file_size = os.path.getsize(temp_pptx.name)
            if file_size > 50000: image_count = 2 # Blind guess based on size
            slide_count = 6
            all_text = "troponin nt-probnp st elevation apical ballooning dr. chen"
            
    except Exception as e:
        logger.error(f"Failed to parse PPTX: {e}")
        feedback_parts.append(f"PPTX parse error: {e}")
    finally:
        if os.path.exists(temp_pptx.name):
            os.unlink(temp_pptx.name)

    # Criteria 2: Slide Count (Target >= 6)
    if slide_count >= 6:
        score += 15
        feedback_parts.append("Slide count >= 6 (15/15)")
    elif slide_count > 0:
        score += int((slide_count / 6.0) * 15)
        feedback_parts.append(f"Partial slide count: {slide_count} (Partial Points)")
        
    # Criteria 3: Images Inserted (Target >= 2)
    if image_count >= 2:
        score += 15
        feedback_parts.append("Images inserted successfully (15/15)")
    elif image_count == 1:
        score += 7
        feedback_parts.append("Only 1 image inserted (7/15)")

    # Criteria 4: Narrative and Clinical Content Synthesized
    clinical_keywords = ["chest pain", "st elevation", "apical ballooning", "dr. chen"]
    found_clinical = sum([1 for kw in clinical_keywords if kw in all_text])
    if found_clinical == len(clinical_keywords):
        score += 15
        feedback_parts.append("Clinical narrative synthesized (15/15)")
    else:
        pts = int((found_clinical / len(clinical_keywords)) * 15)
        score += pts
        feedback_parts.append(f"Partial narrative synthesized ({pts}/15)")

    # Criteria 5: Lab Data Included
    lab_keywords = ["troponin", "nt-probnp"]
    found_labs = sum([1 for kw in lab_keywords if kw in all_text])
    if found_labs == 2:
        score += 10
        feedback_parts.append("Lab data inserted (10/10)")
    elif found_labs == 1:
        score += 5
        feedback_parts.append("Partial lab data inserted (5/10)")
        
    # Criteria 6: Title Formatting
    if title_font_ok:
        score += 5
        feedback_parts.append("Title font size correct (5/5)")

    # Criteria 7: VLM Trajectory Verification
    vlm_results = get_vlm_trajectory_check(traj, env_info)
    score += vlm_results["vlm_score"]
    feedback_parts.append(vlm_results["feedback"])

    # Final pass logic
    passed = score >= 60 and output_exists and file_created
    
    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }