#!/usr/bin/env python3
"""
Verifier for Urban Mobility Vision Zero Pitch task.

Verifies:
1. File exists and was created during the task (anti-gaming).
2. Presentation has at least 5 slides.
3. Content distributed correctly (checking keywords per slide).
4. Images embedded specifically on slides 3 and 5 (index 2 and 4).
5. VLM trajectory verification to ensure UI interaction.
"""

import os
import sys
import json
import logging
import tempfile
import subprocess

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Dynamically import pptx (install if missing on host)
try:
    from pptx import Presentation
except ImportError:
    logger.info("python-pptx not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "python-pptx"])
    from pptx import Presentation


def get_slide_text(slide):
    """Extract all text from a single slide."""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_parts.append(shape.text.lower())
    return " ".join(text_parts)


def has_image(slide):
    """Check if a slide contains an embedded image."""
    for shape in slide.shapes:
        # MSO_SHAPE_TYPE.PICTURE is 13
        if shape.shape_type == 13:
            return True
    return False


def verify_vision_zero_pitch(traj, env_info, task_info):
    """Main verification logic."""
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}
        
    # 1. Read the export JSON from the container
    temp_json = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", temp_json.name)
        with open(temp_json.name, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to read task result: {e}"}
    finally:
        if os.path.exists(temp_json.name):
            os.unlink(temp_json.name)
            
    # Check baseline existence and anti-gaming
    if not result.get("output_exists"):
        return {"passed": False, "score": 0, "feedback": "Failed: Expected output file vision_zero_pitch.pptx not found."}
        
    if not result.get("file_created_during_task"):
        return {"passed": False, "score": 0, "feedback": "Failed: Output file was not modified during the task session (anti-gaming check)."}
        
    if result.get("output_size_bytes", 0) < 10000:
        return {"passed": False, "score": 0, "feedback": "Failed: Output file is suspiciously small (under 10KB)."}

    # 2. Copy the actual PPTX file
    container_pptx_path = result.get("output_path", "/home/ga/Documents/Presentations/Vision_Zero/vision_zero_pitch.pptx")
    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    try:
        copy_from_env(container_pptx_path, temp_pptx.name)
        prs = Presentation(temp_pptx.name)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to parse PPTX document: {e}"}

    # 3. Initialize scoring
    score = 0
    feedback_parts = []
    
    # Check total slide count (10 points)
    num_slides = len(prs.slides)
    if num_slides >= 5:
        score += 10
        feedback_parts.append(f"Slide count check passed ({num_slides} slides).")
    else:
        feedback_parts.append(f"Failed: Insufficient slide count (found {num_slides}, expected at least 5).")

    # Metadata keywords to check per slide
    metadata = task_info.get("metadata", {})
    keywords_by_slide = metadata.get("keywords", {})
    
    # Prevent index errors if fewer slides exist
    s1_text = get_slide_text(prs.slides[0]) if num_slides > 0 else ""
    s2_text = get_slide_text(prs.slides[1]) if num_slides > 1 else ""
    s3_text = get_slide_text(prs.slides[2]) if num_slides > 2 else ""
    s4_text = get_slide_text(prs.slides[3]) if num_slides > 3 else ""
    s5_text = get_slide_text(prs.slides[4]) if num_slides > 4 else ""

    # Check Slide 1 (10 pts)
    k1 = keywords_by_slide.get("slide_1", ["hoboken", "vision zero"])
    if all(k in s1_text for k in k1):
        score += 10
        feedback_parts.append("Slide 1 content check passed.")
    else:
        feedback_parts.append("Slide 1 missing required keywords.")

    # Check Slide 2 (10 pts)
    k2 = keywords_by_slide.get("slide_2", ["vulnerable", "fatalities"])
    if all(k in s2_text for k in k2):
        score += 10
        feedback_parts.append("Slide 2 content check passed.")
    else:
        feedback_parts.append("Slide 2 missing required keywords.")

    # Check Slide 3 (15 pts content, 15 pts image)
    k3 = keywords_by_slide.get("slide_3", ["daylighting", "infrastructure"])
    if all(k in s3_text for k in k3):
        score += 15
        feedback_parts.append("Slide 3 content check passed.")
    else:
        feedback_parts.append("Slide 3 missing required keywords.")
        
    if num_slides > 2 and has_image(prs.slides[2]):
        score += 15
        feedback_parts.append("Slide 3 image check passed.")
    else:
        feedback_parts.append("Slide 3 missing embedded image.")

    # Check Slide 4 (10 pts)
    k4 = keywords_by_slide.get("slide_4", ["implementation", "timeline"])
    if all(k in s4_text for k in k4):
        score += 10
        feedback_parts.append("Slide 4 content check passed.")
    else:
        feedback_parts.append("Slide 4 missing required keywords.")

    # Check Slide 5 (15 pts content, 15 pts image)
    k5 = keywords_by_slide.get("slide_5", ["adopt resolution", "commit"])
    if all(k in s5_text for k in k5):
        score += 15
        feedback_parts.append("Slide 5 content check passed.")
    else:
        feedback_parts.append("Slide 5 missing required keywords.")

    if num_slides > 4 and has_image(prs.slides[4]):
        score += 15
        feedback_parts.append("Slide 5 image check passed.")
    else:
        feedback_parts.append("Slide 5 missing embedded image.")

    # Cleanup temp file
    if os.path.exists(temp_pptx.name):
        os.unlink(temp_pptx.name)

    # 4. Optional Trajectory VLM Check
    # This checks that ONLYOFFICE was actively used, verifying the workflow.
    try:
        from vlm_utils import sample_trajectory_frames, query_vlm
        
        frames = sample_trajectory_frames(traj, n=4)
        if frames:
            prompt = """Look at these screenshots from a desktop environment.
            Determine if the user was actively using presentation software (like ONLYOFFICE or PowerPoint) 
            to create or edit a slide deck.
            
            Respond in JSON: {"presentation_software_used": true/false}"""
            
            vlm_response = query_vlm(images=frames, prompt=prompt)
            if vlm_response.get("success") and vlm_response.get("parsed", {}).get("presentation_software_used"):
                feedback_parts.append("VLM verified presentation software usage.")
            else:
                feedback_parts.append("VLM could not confirm presentation software usage (possible penalty).")
                # Deduct slight points if clear UI interaction isn't verified
                score = max(0, score - 10)
    except Exception as e:
        logger.warning(f"VLM verification skipped or failed: {e}")

    # Pass condition: must reach 70 points
    passed = score >= 70

    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }