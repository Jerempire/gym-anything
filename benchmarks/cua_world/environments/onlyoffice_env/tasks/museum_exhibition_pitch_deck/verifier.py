#!/usr/bin/env python3
"""
Verifier for Museum Exhibition Pitch Deck task.

Verifies:
1. File exists and was created during the task (anti-gaming).
2. Presentation has at least 8 slides.
3. Slide 1 contains the required speaker notes ("Miller Foundation").
4. Presentation contains at least 5 images (4 artifacts + 1 floorplan).
5. Presentation contains a table with budget figures ("150,000", "Fabrication").
6. VLM trajectory verification to ensure agent actually assembled the deck.
"""

import json
import os
import sys
import tempfile
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

sys.path.insert(0, str(os.path.abspath(os.path.join(os.path.dirname(__file__), '../../'))))
from gym_anything.vlm import sample_trajectory_frames, get_final_screenshot, query_vlm

# We need python-pptx to parse the output file.
# It is installed in the onlyoffice_env per environment scripts.
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx is not installed. Will attempt to install.")

def ensure_pptx():
    global PPTX_AVAILABLE
    if not PPTX_AVAILABLE:
        try:
            import subprocess
            subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "python-pptx"])
            global Presentation, MSO_SHAPE_TYPE
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            PPTX_AVAILABLE = True
        except Exception as e:
            logger.error(f"Failed to install python-pptx: {e}")
            return False
    return True

def extract_all_text(shape):
    """Recursively extract text from a shape and its children."""
    text = ""
    if shape.has_text_frame:
        text += shape.text + " "
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text_frame.text + " "
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            text += extract_all_text(child)
    return text

def verify_presentation_assembly(traj, env_info, task_info):
    """
    Programmatic and VLM-based verification for the presentation assembly task.
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    if not ensure_pptx():
        return {"passed": False, "score": 0, "feedback": "python-pptx library unavailable"}

    score = 0
    feedback_parts = []
    
    # Extract metadata
    metadata = task_info.get('metadata', {})
    expected_path = metadata.get('expected_output_path', '/home/ga/Documents/Presentations/exhibition_proposal.pptx')
    
    # Copy task_result.json
    temp_res = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", temp_res.name)
        with open(temp_res.name, 'r') as f:
            result_data = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to read result: {e}"}
    finally:
        if os.path.exists(temp_res.name):
            os.unlink(temp_res.name)
            
    # Check if file exists and was created
    output_exists = result_data.get('output_exists', False)
    file_created = result_data.get('file_created_during_task', False)
    
    if not output_exists:
        return {
            "passed": False,
            "score": 0,
            "feedback": "exhibition_proposal.pptx was not found. Agent failed to save the file."
        }
        
    if file_created:
        score += 15
        feedback_parts.append("File created during session")
    else:
        feedback_parts.append("File existed before task (Warning: possible gaming)")

    # Copy the PPTX file for parsing
    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    try:
        copy_from_env(expected_path, temp_pptx.name)
        prs = Presentation(temp_pptx.name)
        
        # 1. Check Slide Count (15 points)
        slide_count = len(prs.slides)
        if slide_count >= 8:
            score += 15
            feedback_parts.append(f"Slide count acceptable ({slide_count})")
        else:
            score += int((slide_count / 8) * 15)
            feedback_parts.append(f"Insufficient slides ({slide_count}/8)")

        # Variables for parsing
        has_miller_note = False
        image_count = 0
        has_budget_table = False
        has_budget_text = False
        
        # 2. Iterate through slides to verify content
        for i, slide in enumerate(prs.slides):
            # Check speaker notes on Slide 1
            if i == 0 and slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.lower()
                if "miller" in notes_text and "foundation" in notes_text:
                    has_miller_note = True

            for shape in slide.shapes:
                # Count Images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, 'image'):
                    image_count += 1
                    
                # Check for Budget Table
                if shape.has_table:
                    table_text = extract_all_text(shape).lower()
                    if "150,000" in table_text and "fabrication" in table_text:
                        has_budget_table = True
                        
                # Just in case they used a text box instead of a table
                if shape.has_text_frame:
                    text = shape.text.lower()
                    if "150,000" in text and "fabrication" in text:
                        has_budget_text = True

        # Evaluate Notes (10 points)
        if has_miller_note:
            score += 10
            feedback_parts.append("Speaker notes added correctly")
        else:
            feedback_parts.append("Missing required speaker notes")
            
        # Evaluate Images (25 points)
        if image_count >= 5:
            score += 25
            feedback_parts.append(f"Found {image_count} images (Pass)")
        elif image_count > 0:
            score += int((image_count / 5) * 25)
            feedback_parts.append(f"Found {image_count} images (Partial)")
        else:
            feedback_parts.append("No images found in presentation")
            
        # Evaluate Budget Table (20 points)
        if has_budget_table:
            score += 20
            feedback_parts.append("Budget table found and correct")
        elif has_budget_text:
            score += 10
            feedback_parts.append("Budget data found but not in table format")
        else:
            feedback_parts.append("Budget data missing")

    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        feedback_parts.append(f"Error parsing PPTX file: {str(e)}")
    finally:
        if os.path.exists(temp_pptx.name):
            os.unlink(temp_pptx.name)

    # VLM Trajectory Verification (15 points)
    # This prevents the agent from somehow cheating the file creation without using the UI
    try:
        frames = sample_trajectory_frames(traj, n=4)
        final_img = get_final_screenshot(traj)
        
        prompt = """
        Review these screenshots of a computer agent using ONLYOFFICE Presentation Editor.
        Task: Assemble a presentation pitch deck with images, text, and tables.
        
        Did the agent visibly interact with the presentation software, insert images, and add text across multiple slides?
        Respond in JSON:
        {
            "interacted_with_software": true/false,
            "inserted_images": true/false,
            "added_text": true/false,
            "confidence": "high/medium/low",
            "reasoning": "Brief explanation"
        }
        """
        
        vlm_resp = query_vlm(images=frames + [final_img], prompt=prompt)
        if vlm_resp.get("success"):
            parsed = vlm_resp.get("parsed", {})
            vlm_score = 0
            if parsed.get("interacted_with_software"): vlm_score += 5
            if parsed.get("inserted_images"): vlm_score += 5
            if parsed.get("added_text"): vlm_score += 5
            
            score += vlm_score
            feedback_parts.append(f"VLM Visual Verification: {vlm_score}/15 pts")
        else:
            feedback_parts.append("VLM Verification failed or unavailable")
            
    except Exception as e:
        logger.error(f"VLM verification error: {e}")
        feedback_parts.append("VLM Trajectory check skipped")

    # Determine passing status
    # Must have >= 70 points AND must have actually created a file with images
    key_criteria_met = file_created and (image_count >= 2) and (slide_count >= 6)
    passed = (score >= 70) and key_criteria_met

    return {
        "passed": passed,
        "score": min(score, 100),
        "feedback": " | ".join(feedback_parts)
    }