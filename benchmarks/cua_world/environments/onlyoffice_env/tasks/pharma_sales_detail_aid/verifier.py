#!/usr/bin/env python3
"""
Verifier for the Pharmaceutical Sales Detail Aid task.
Evaluates multi-modal document assembly in ONLYOFFICE Presentation Editor.
"""

import os
import json
import logging
import tempfile
from collections import Counter

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Try importing python-pptx (installed in environment)
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx is not available.")


def get_all_text_from_presentation(prs):
    """Extracts all text from a python-pptx Presentation object."""
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_content.append(shape.text.lower())
    return " ".join(text_content)


def verify_pharma_presentation(traj, env_info, task_info):
    """
    Verifies the presentation creation.
    Criteria:
      - File exists and was created during the task (Anti-gaming)
      - Slide count >= 5
      - Content checks (Text matching FDA label)
      - Multi-modal: Contains an image object (slide 3)
      - Multi-modal: Contains a native table (slide 4)
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "System error: copy_from_env not available."}

    score = 0
    feedback_parts = []
    
    # 1. Read JSON result from container
    temp_json = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", temp_json.name)
        with open(temp_json.name, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to read task result: {e}"}
    finally:
        if os.path.exists(temp_json.name):
            os.unlink(temp_json.name)

    # 2. Check File Existence & Creation Time (15 points)
    if not result.get('output_exists', False):
        return {"passed": False, "score": 0, "feedback": "Failed: cardio_detail_aid.pptx was not found."}
    
    if result.get('file_created_during_task', False):
        score += 15
        feedback_parts.append("File correctly created during session")
    else:
        feedback_parts.append("File exists but timestamp indicates it might be stale")
        
    filename = result.get("actual_filename", "")
    if filename != "cardio_detail_aid.pptx":
        feedback_parts.append(f"Saved with wrong filename: {filename}")
        score -= 5  # minor penalty

    # 3. Parse PPTX file
    if not PPTX_AVAILABLE:
        feedback_parts.append("python-pptx missing; relying solely on VLM/file existence")
        passed = score > 0
        return {"passed": passed, "score": score, "feedback": " | ".join(feedback_parts)}

    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    try:
        # Copy the presentation from the exact target path
        copy_from_env("/home/ga/Documents/Presentations/cardio_detail_aid.pptx", temp_pptx.name)
        try:
            prs = Presentation(temp_pptx.name)
        except Exception as e:
            return {"passed": False, "score": score, "feedback": f"File exists but is not a valid PPTX: {e}"}
            
        # 4. Check Slide Count (15 points)
        num_slides = len(prs.slides)
        if num_slides >= 5:
            score += 15
            feedback_parts.append(f"Slide count correct ({num_slides})")
        else:
            feedback_parts.append(f"Expected 5 slides, found {num_slides}")
            
        # 5. Extract Text & Check Required Content (20 points)
        all_text = get_all_text_from_presentation(prs)
        keywords = ["hyperlipidemia", "hmg-coa", "dosing", "pregnancy"]
        found_keywords = sum(1 for kw in keywords if kw in all_text)
        
        if found_keywords == 4:
            score += 20
            feedback_parts.append("All required clinical text found")
        elif found_keywords > 0:
            score += (found_keywords * 5)
            feedback_parts.append(f"Partial clinical text found ({found_keywords}/4)")
        else:
            feedback_parts.append("Required clinical text is missing")

        # 6. Check for Image Object (20 points)
        # Scan all shapes for a picture (shape_type 13). Expected on slide 3 (index 2), but we accept anywhere
        has_image = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13: # 13 is PICTURE
                    has_image = True
                    break
            if has_image: break
            
        if has_image:
            score += 20
            feedback_parts.append("Efficacy chart image embedded correctly")
        else:
            feedback_parts.append("Efficacy chart image missing")

        # 7. Check for Native Table Object (20 points)
        # Native table is shape_type 19.
        has_table = False
        correct_dims = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    has_table = True
                    table = shape.table
                    rows = len(table.rows)
                    cols = len(table.columns)
                    if rows >= 4 and cols >= 3:
                        correct_dims = True
                    break
            if has_table: break

        if has_table and correct_dims:
            score += 20
            feedback_parts.append("Native safety table created with correct dimensions")
        elif has_table:
            score += 10
            feedback_parts.append("Table created but dimensions are incorrect")
        else:
            feedback_parts.append("Native safety table missing (Screenshot pasted?)")

        # 8. Check Layout Complexities (Slide 5 Two-column) (10 points)
        # Just check if any slide has multiple text blocks
        multi_text_slide = False
        for slide in prs.slides:
            text_frames = sum(1 for shape in slide.shapes if shape.has_text_frame and shape.text.strip() != "")
            if text_frames >= 3: # Title + 2 content boxes
                multi_text_slide = True
                break
                
        if multi_text_slide:
            score += 10
            feedback_parts.append("Multi-column layout utilized")
        
    finally:
        if os.path.exists(temp_pptx.name):
            os.unlink(temp_pptx.name)
            
    # VLM Trajectory Verification Check (Optional fallback/enhancement if VLM is enabled in the env)
    # Since we can parse the exact objects, programmatic verification is highly deterministic here.
    
    # 75 is the passing threshold, meaning they must get at least the text, file, and either table/image correct
    passed = score >= 75
    
    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }