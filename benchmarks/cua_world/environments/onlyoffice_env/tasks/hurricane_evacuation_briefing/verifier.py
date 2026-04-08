#!/usr/bin/env python3
"""
Verifier for Hurricane Emergency Evacuation Briefing task.

1. File Creation & Anti-Gaming: File exists and was created during task.
2. Presentation Parsing (via python-pptx):
   - Slide Count (6)
   - Slide 1: Title & Subtitle
   - Slide 2: Storm Stats
   - Slide 3 & 4: Images Inserted
   - Slide 4: Surge heights stated
   - Slide 5 & 6: Evacuation Zones & Shelters
   - Slide 6: Speaker Notes containing threshold condition
3. VLM Trajectory Verification: Ensures ONLYOFFICE UI was used to construct the deck.
"""

import os
import json
import logging
import tempfile
from typing import Dict, Any

# Gym-anything VLM utilities
from gym_anything.vlm import sample_trajectory_frames, query_vlm

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Try to import python-pptx (available in environment)
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. Structural verification will fail.")


def extract_text_from_slide(slide) -> str:
    """Extract all text from all shapes in a slide."""
    text_content = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_content.append(shape.text.lower())
    return " ".join(text_content)


def count_images_on_slide(slide) -> int:
    """Count picture shapes (shape_type == 13) on a slide."""
    img_count = 0
    for shape in slide.shapes:
        # 13 is MSO_SHAPE_TYPE.PICTURE
        if getattr(shape, "shape_type", None) == 13 or (hasattr(shape, "name") and "Picture" in shape.name):
            img_count += 1
    return img_count


def verify_hurricane_briefing(traj: Dict[str, Any], env_info: Dict[str, Any], task_info: Dict[str, Any]) -> Dict[str, Any]:
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    metadata = task_info.get('metadata', {})
    expected_output = metadata.get('expected_output_path')
    
    score = 0
    feedback_parts = []
    
    # 1. Read JSON result
    temp_json = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", temp_json.name)
        with open(temp_json.name, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to read result: {e}"}
    finally:
        if os.path.exists(temp_json.name):
            os.unlink(temp_json.name)

    # File Creation Check (10 points)
    if result.get("output_exists") and result.get("file_created_during_task"):
        score += 10
        feedback_parts.append("File created successfully")
    else:
        feedback_parts.append("Output file missing or not created during task window")
        return {"passed": False, "score": 0, "feedback": " | ".join(feedback_parts)}

    # 2. PPTX Structural Analysis
    if not PPTX_AVAILABLE:
        return {"passed": False, "score": 0, "feedback": "python-pptx library missing"}

    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    try:
        copy_from_env(expected_output, temp_pptx.name)
        prs = Presentation(temp_pptx.name)
        
        slides = prs.slides
        slide_count = len(slides)
        
        # Slide Count (10 points)
        if slide_count == metadata.get("expected_slide_count", 6):
            score += 10
            feedback_parts.append("Exactly 6 slides found")
        else:
            feedback_parts.append(f"Expected 6 slides, found {slide_count}")
            # Pad slides list to avoid index errors if agent made too few slides
            while len(slides) < 6:
                slides.append(None)

        slide_texts = [extract_text_from_slide(s) if s else "" for s in slides]

        # Title Content - Slide 1 (15 points)
        if all(term in slide_texts[0] for term in metadata.get("expected_title_terms", [])):
            score += 15
            feedback_parts.append("Title/Subtitle correct")
        else:
            feedback_parts.append("Missing required terms on Slide 1")

        # Status Data - Slide 2 (15 points)
        if all(term in slide_texts[1] for term in metadata.get("expected_status_terms", [])):
            score += 15
            feedback_parts.append("Storm stats accurately extracted")
        else:
            feedback_parts.append("Missing required storm stats on Slide 2")

        # Graphical Integration - Slides 3 & 4 (20 points)
        s3_images = count_images_on_slide(slides[2]) if slides[2] else 0
        s4_images = count_images_on_slide(slides[3]) if slides[3] else 0
        if s3_images > 0 and s4_images > 0:
            score += 20
            feedback_parts.append("Graphics successfully inserted")
        elif s3_images > 0 or s4_images > 0:
            score += 10
            feedback_parts.append("Partial graphics inserted")
        else:
            feedback_parts.append("No graphics found on Slides 3 or 4")

        # Surge Data - Slide 4 (10 points)
        if all(term in slide_texts[3] for term in metadata.get("expected_surge_terms", [])):
            score += 10
            feedback_parts.append("Surge height correctly stated")
        else:
            feedback_parts.append("Missing surge height text on Slide 4")

        # Operations Data - Slides 5 & 6 (10 points)
        s5_evac = all(term in slide_texts[4] for term in metadata.get("expected_evac_terms", []))
        s6_shelter = all(term in slide_texts[5] for term in metadata.get("expected_shelter_terms", []))
        if s5_evac and s6_shelter:
            score += 10
            feedback_parts.append("Evacuation & Shelter info correct")
        else:
            feedback_parts.append("Missing operations data on Slides 5/6")

        # Speaker Notes - Slide 6 (10 points)
        notes = ""
        if slides[5] and slides[5].has_notes_slide:
            notes = slides[5].notes_slide.notes_text_frame.text.lower()
            
        if metadata.get("expected_speaker_note") in notes:
            score += 10
            feedback_parts.append("Speaker note successfully added")
        else:
            feedback_parts.append("Speaker note missing or incorrect on Slide 6")

    except Exception as e:
        feedback_parts.append(f"PPTX parsing error: {str(e)}")
    finally:
        if os.path.exists(temp_pptx.name):
            os.unlink(temp_pptx.name)

    # 3. VLM Trajectory Verification
    # We ask the VLM to confirm the agent actually used the ONLYOFFICE Presentation UI
    frames = sample_trajectory_frames(traj, n=4)
    vlm_prompt = """Look at these screenshots taken during the agent's workflow.
Did the agent use a presentation software (like ONLYOFFICE) to create slides, edit text, and insert pictures?
We want to verify this was done via the graphical UI, not just via a background script.
Answer in JSON: {"used_ui": true/false, "reason": "brief explanation"}"""
    
    vlm_res = query_vlm(images=frames, prompt=vlm_prompt)
    used_ui = False
    if vlm_res and vlm_res.get("success") and vlm_res.get("parsed"):
        used_ui = vlm_res["parsed"].get("used_ui", False)
        if used_ui:
            feedback_parts.append("VLM confirmed UI usage")
        else:
            feedback_parts.append("VLM did not detect presentation UI usage (possible script bypass)")
            score -= 30  # Penalty for gaming the UI requirement

    # Pass logic: Needs 70 points, must have created the file (10), must have graphics (20), must have used UI
    key_criteria_met = (result.get("file_created_during_task") and 
                       (s3_images > 0 or s4_images > 0) and 
                       used_ui)
    
    passed = (score >= 70) and key_criteria_met

    return {
        "passed": passed,
        "score": max(0, score),
        "feedback": " | ".join(feedback_parts)
    }