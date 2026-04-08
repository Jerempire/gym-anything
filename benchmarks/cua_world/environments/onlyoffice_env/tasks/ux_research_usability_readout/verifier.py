#!/usr/bin/env python3
import json
import os
import tempfile
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Fallback imports for python-pptx (installed in env)
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available.")

def verify_ux_readout(traj, env_info, task_info):
    """
    Verify the UX Readout presentation task.
    
    Scores (100 pts max):
    1. File exists and created during task (15 pts) - ANTI-GAMING
    2. Slide count >= 5 (10 pts)
    3. Mentions methodology/sample size "30" (10 pts)
    4. Correct SUS Score "68.5" (20 pts)
    5. Correct Success Rate "82.2%" (20 pts)
    6. Contains at least two quotes (15 pts)
    7. Contains prototype image (10 pts)
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    metadata = task_info.get('metadata', {})
    expected_output_path = metadata.get('expected_output_path', '/home/ga/Documents/Presentations/ux_research_readout.pptx')
    quotes_path = metadata.get('quotes_file', '/home/ga/Documents/Presentations/user_quotes.txt')
    
    score = 0
    feedback_parts = []
    
    # 1. Read task execution result
    temp_result = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    result_data = {}
    try:
        copy_from_env("/tmp/task_result.json", temp_result.name)
        with open(temp_result.name, 'r') as f:
            result_data = json.load(f)
    except Exception as e:
        logger.error(f"Failed to read /tmp/task_result.json: {e}")
    finally:
        if os.path.exists(temp_result.name):
            os.unlink(temp_result.name)
            
    # CRITERION 1: File Existence & Anti-Gaming
    output_exists = result_data.get('output_exists', False)
    created_during = result_data.get('file_created_during_task', False)
    
    if not output_exists:
        return {"passed": False, "score": 0, "feedback": "PPTX output file was not found."}
    
    if created_during:
        score += 15
        feedback_parts.append("File created/modified during task (+15)")
    else:
        feedback_parts.append("Warning: File not modified during task timeframe.")

    if not PPTX_AVAILABLE:
        # If we can't parse PPTX, rely entirely on VLM trajectory verification
        return fallback_vlm_verification(traj, env_info, score, feedback_parts)

    # 2. Parse PPTX
    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    temp_quotes = tempfile.NamedTemporaryFile(delete=False, suffix='.txt')
    try:
        copy_from_env(expected_output_path, temp_pptx.name)
        copy_from_env(quotes_path, temp_quotes.name)
        
        prs = Presentation(temp_pptx.name)
        
        # Extract text and images
        all_text = ""
        has_image = False
        num_slides = len(prs.slides)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + " \n"
                if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                    has_image = True
                    
        all_text_lower = all_text.lower()
        
        # Read quotes
        with open(temp_quotes.name, 'r') as f:
            raw_quotes = f.read()
        # Extract meaningful snippets to search for
        quote_snippets = [
            "payment went through",
            "auto-fill feature",
            "find where to enter my discount code",
            "text color on the 'submit order'",
            "guest checkout is a lifesaver"
        ]
        
        # CRITERION 2: Slide Count >= 5
        if num_slides >= 5:
            score += 10
            feedback_parts.append(f"Slide count {num_slides} >= 5 (+10)")
        else:
            feedback_parts.append(f"Slide count is {num_slides}, expected >= 5.")

        # CRITERION 3: Sample Size
        if "30" in all_text:
            score += 10
            feedback_parts.append("Methodology mentions sample size 30 (+10)")
        else:
            feedback_parts.append("Methodology sample size '30' not found.")

        # CRITERION 4: SUS Score
        # Accept exactly 68.5 or rounded 68/69
        if "68.5" in all_text or "68" in all_text or "69" in all_text:
            score += 20
            feedback_parts.append("Correct SUS Score calculated (+20)")
        else:
            feedback_parts.append("Calculated SUS score (68.5) not found.")

        # CRITERION 5: Success Rate
        # Accept exactly 82.2% or rounded 82%
        if "82.2" in all_text or "82" in all_text or "0.82" in all_text:
            score += 20
            feedback_parts.append("Correct Success Rate calculated (+20)")
        else:
            feedback_parts.append("Calculated Success Rate (82.2%) not found.")

        # CRITERION 6: Quotes
        found_quotes = 0
        for snippet in quote_snippets:
            if snippet.lower() in all_text_lower:
                found_quotes += 1
                
        if found_quotes >= 2:
            score += 15
            feedback_parts.append(f"Found {found_quotes} verbatim quotes (+15)")
        elif found_quotes == 1:
            score += 7
            feedback_parts.append(f"Found only 1 verbatim quote (+7)")
        else:
            feedback_parts.append("Verbatim quotes not found in presentation.")

        # CRITERION 7: Image inserted
        if has_image:
            score += 10
            feedback_parts.append("Prototype image found in presentation (+10)")
        else:
            feedback_parts.append("No images found in the presentation.")

    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        feedback_parts.append(f"Error parsing PPTX file: {str(e)}")
    finally:
        if os.path.exists(temp_pptx.name):
            os.unlink(temp_pptx.name)
        if os.path.exists(temp_quotes.name):
            os.unlink(temp_quotes.name)

    # Evaluate final
    passed = score >= 70
    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }

def fallback_vlm_verification(traj, env_info, current_score, feedback_parts):
    """
    Fallback verification using VLM if python-pptx is completely unavailable.
    """
    try:
        from gym_anything.vlm import sample_trajectory_frames, get_final_screenshot, query_vlm
        
        frames = sample_trajectory_frames(traj, n=3)
        final = get_final_screenshot(traj)
        
        prompt = """
        Review these screenshots of an agent creating an executive readout presentation in ONLYOFFICE.
        Check for the following:
        1. Is there a presentation with multiple slides visible?
        2. Can you see calculated metrics like '68.5' (SUS) or '82.2%' (Success Rate)?
        3. Can you see user quotes pasted into the slides?
        4. Is the prototype image (a wireframe or mockup of Express Checkout) inserted into a slide?
        
        Respond with JSON containing boolean keys: "has_multiple_slides", "has_correct_metrics", "has_quotes", "has_image"
        """
        
        vlm_result = query_vlm(images=frames + [final], prompt=prompt)
        parsed = vlm_result.get("parsed", {})
        
        vlm_score = 0
        if parsed.get("has_multiple_slides"): vlm_score += 10
        if parsed.get("has_correct_metrics"): vlm_score += 40
        if parsed.get("has_quotes"): vlm_score += 15
        if parsed.get("has_image"): vlm_score += 10
        
        total_score = current_score + vlm_score
        feedback_parts.append(f"VLM Verification Score: {vlm_score}")
        
        return {
            "passed": total_score >= 70,
            "score": total_score,
            "feedback": " | ".join(feedback_parts)
        }
    except Exception as e:
        return {
            "passed": False,
            "score": current_score,
            "feedback": " | ".join(feedback_parts) + f" | VLM fallback failed: {e}"
        }